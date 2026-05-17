"""Generate trybe-deck.pptx from slide content."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
import copy
from lxml import etree

# ── Colours ───────────────────────────────────────────────────────────────────
PINK    = RGBColor(0xE6, 0x00, 0x7E)
INK     = RGBColor(0x1A, 0x09, 0x18)
INK2    = RGBColor(0x4A, 0x2F, 0x47)
MUTED   = RGBColor(0x8A, 0x75, 0x85)
BG      = RGBColor(0xF8, 0xF4, 0xF0)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
RULE    = RGBColor(0xE0, 0xD5, 0xC8)
PINK_T  = RGBColor(0xFC, 0xE8, 0xF2)
GLASS   = RGBColor(0xFF, 0xFF, 0xFF)

# ── Slide dimensions (16:9 widescreen) ────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # blank

# ── Helpers ───────────────────────────────────────────────────────────────────

def add_slide():
    s = prs.slides.add_slide(blank_layout)
    bg_fill(s)
    return s

def bg_fill(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG

def txb(slide, text, x, y, w, h,
        size=14, bold=False, color=INK, align=PP_ALIGN.LEFT,
        font="Calibri", wrap=True, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return tb

def label_txt(slide, text, x, y, w=Inches(10), color=PINK, size=9):
    return txb(slide, text, x, y, w, Pt(20), size=size, color=color,
               font="Calibri", bold=False)

def heading(slide, text, x, y, w, size=36, color=INK):
    return txb(slide, text, x, y, w, Inches(2), size=size, bold=False,
               color=color, font="Calibri")

def rect(slide, x, y, w, h, fill=WHITE, alpha=None, line_color=RULE, line_width=Pt(0.5)):
    shp = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line_color
    shp.line.width = line_width
    return shp

def add_text_to_shape(shp, text, size=12, bold=False, color=INK, align=PP_ALIGN.LEFT, font="Calibri"):
    tf = shp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font

def chrome_top(slide, right_label):
    """Top chrome bar: TRYBE brand left, label right."""
    txb(slide, "▶  TRYBE", Inches(0.7), Inches(0.45), Inches(3), Pt(18),
        size=8, bold=True, color=INK, font="Calibri")
    txb(slide, right_label, Inches(8), Inches(0.45), Inches(5), Pt(18),
        size=8, color=MUTED, align=PP_ALIGN.RIGHT, font="Calibri")

def chrome_bot(slide, left="", right=""):
    if left:
        txb(slide, left, Inches(0.7), Inches(7.1), Inches(6), Pt(18),
            size=8, color=MUTED, font="Calibri")
    if right:
        txb(slide, right, Inches(7), Inches(7.1), Inches(6), Pt(18),
            size=8, color=MUTED, align=PP_ALIGN.RIGHT, font="Calibri")

def glass_box(slide, x, y, w, h):
    """White semi-transparent box with rule border."""
    r = rect(slide, x, y, w, h, fill=RGBColor(0xFF,0xFF,0xFF), line_color=RULE)
    return r

def bullet_para(tf, text, size=12, color=INK2, first=False, indent=0):
    if first:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = ""
    p.space_before = Pt(4)
    run = p.add_run()
    run.text = f"▸  {text}" if not indent else f"    {text}"
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return p

def section_title_block(slide, label, title, subtitle=None,
                         lx=Inches(0.9), ly=Inches(0.85),
                         lw=Inches(11.5)):
    """Label + h3 + optional subtitle."""
    label_txt(slide, label, lx, ly, lw)
    hy = ly + Inches(0.32)
    txb(slide, title, lx, hy, lw, Inches(1.2),
        size=34, color=INK, font="Calibri")
    if subtitle:
        sy = hy + Inches(0.9)
        txb(slide, subtitle, lx, sy, lw, Inches(0.7),
            size=13, color=INK2, font="Calibri")

# ── Slide 1 — Cover ────────────────────────────────────────────────────────────
s1 = add_slide()

# Pink accent bar top
r = rect(s1, Inches(0), Inches(0), Inches(13.33), Inches(0.08), fill=PINK, line_color=PINK)

# Brand
txb(s1, "▶  TRYBE × CANDIDATE", Inches(0.9), Inches(0.55), Inches(6), Pt(20),
    size=10, bold=True, color=INK, font="Calibri")
txb(s1, "Head of People · May 2026", Inches(7), Inches(0.55), Inches(5.5), Pt(20),
    size=9, color=MUTED, align=PP_ALIGN.RIGHT, font="Calibri")

# Label
label_txt(s1, "— People strategy as operating system", Inches(0.9), Inches(2.0), color=PINK)

# Big headline
txb(s1, "People strategy as the operating system for the pivot.",
    Inches(0.9), Inches(2.45), Inches(11.2), Inches(2.0),
    size=42, bold=False, color=INK, font="Calibri")

# Subtitle
txb(s1, "From spa software to experience platform — designing the People function for the journey.",
    Inches(0.9), Inches(4.6), Inches(10.0), Inches(0.6),
    size=14, color=INK2, font="Calibri")

# Meta grid
meta = [
    ("Presenter", "Stephanie Fadahunsi · Head of People candidate"),
    ("For",       "Trybe · Series A · Five Elms Capital · London"),
    ("Date",      "May 2026"),
]
my = Inches(5.5)
for k, v in meta:
    txb(s1, k.upper(), Inches(0.9), my, Inches(1.4), Pt(20),
        size=8, color=MUTED, font="Calibri")
    txb(s1, v, Inches(2.4), my, Inches(9.0), Pt(20),
        size=12, color=INK2, font="Calibri")
    my += Inches(0.32)

# ── Slide 2 — Today and becoming ───────────────────────────────────────────────
s2 = add_slide()
chrome_top(s2, "02 · Today and becoming")

label_txt(s2, "How I see Trybe", Inches(0.9), Inches(0.85))
txb(s2, "Today and becoming.",
    Inches(0.9), Inches(1.2), Inches(11.5), Inches(0.8),
    size=34, color=INK, font="Calibri")
txb(s2, "The People function has to serve both — operating for the business today, and building for the company Five Elms is underwriting.",
    Inches(0.9), Inches(2.05), Inches(11.5), Inches(0.55),
    size=12, color=INK2, font="Calibri")

# Table
bx = Inches(0.9)
by = Inches(2.7)
bw = Inches(11.5)
bh = Inches(4.5)
glass_box(s2, bx, by, bw, bh)

headers = ["", "Today", "Becoming"]
col_ws  = [Inches(1.6), Inches(4.5), Inches(5.4)]

rows = [
    # (type, cells)
    ("sect", ["VALUE CREATION", "", ""]),
    ("data", ["Category",    "Spa software",
              "Experience platform for hospitality"]),
    ("data", ["What it sells",
              "Spa operations: bookings, scheduling, payments, memberships",
              "The system underneath everything in the hotel that isn't the bedroom"]),
    ("data", ["Who buys",
              "The Spa Director, with a spa budget",
              "The GM and CFO, with a hotel-wide budget"]),
    ("sect", ["VALUE CAPTURE", "", ""]),
    ("data", ["Revenue",
              "Modular subscription: Basic, Pro, Enterprise + small payments take rate",
              "Multi-module subscriptions, payments take rate, AI features; growing automatically as customers deepen use"]),
    ("data", ["Lock-in",
              "Once installed, leaving means reconfiguring every treatment and moving years of data.",
              "The same friction, spanning far more of the hotel. Replacing Trybe means replacing the software underneath most of the hotel's revenue."]),
    ("data", ["Moat",
              "Works for the spa. Doesn't extend further.",
              "Structural. Payments lock-in deepens with volume. AI compounds with data."]),
    ("data", ["Commercial",
              "Under-priced relative to value created",
              "Captures the share of value that category-defining platforms earn"]),
]

# header row
row_h = Pt(20)
hx = bx + Inches(0.05)
hy = by + Inches(0.06)
cx = hx
for i, (hdr, cw) in enumerate(zip(headers, col_ws)):
    c = MUTED if i == 0 else (MUTED if i == 1 else PINK)
    txb(s2, hdr.upper(), cx, hy, cw - Inches(0.1), Pt(18),
        size=8, color=c, font="Calibri")
    cx += cw

row_y = hy + Inches(0.28)
row_height = Inches(0.44)

for rtype, cells in rows:
    if rtype == "sect":
        bg = rect(s2, bx, row_y, bw, Inches(0.28), fill=PINK_T, line_color=RULE, line_width=Pt(0.25))
        txb(s2, cells[0], bx + Inches(0.15), row_y + Inches(0.04), bw - Inches(0.2), Inches(0.25),
            size=8, color=PINK, font="Calibri")
        row_y += Inches(0.28)
    else:
        cx2 = bx + Inches(0.05)
        colors = [MUTED, INK2, INK]
        bolds  = [False, False, False]
        for i, (cell, cw) in enumerate(zip(cells, col_ws)):
            txb(s2, cell, cx2 + Inches(0.05), row_y + Inches(0.04),
                cw - Inches(0.15), row_height - Inches(0.06),
                size=10, color=colors[i], bold=bolds[i], font="Calibri")
            cx2 += cw
        # bottom border line
        ln = rect(s2, bx, row_y + row_height, bw, Emu(5000), fill=RULE, line_color=RULE, line_width=Emu(5000))
        row_y += row_height

# ── Slide 3 — What decides ─────────────────────────────────────────────────────
s3 = add_slide()
chrome_top(s3, "03 · Strategic priorities — ranked")

label_txt(s3, "Ranked strategic analysis — not balanced", Inches(0.9), Inches(0.85))
txb(s3, "What actually decides the pivot.",
    Inches(0.9), Inches(1.2), Inches(11.5), Inches(0.8),
    size=34, color=INK, font="Calibri")

# Two determinants label
txb(s3, "The two determinants — where 80% of the outcome lives",
    Inches(0.9), Inches(2.1), Inches(11.5), Pt(18), size=9, color=INK2, font="Calibri")

# Card 01
c1 = glass_box(s3, Inches(0.9), Inches(2.4), Inches(5.65), Inches(1.9))
txb(s3, "01", Inches(1.0), Inches(2.5), Inches(1.5), Pt(30),
    size=22, color=PINK, font="Calibri")
txb(s3, "[THREAT]", Inches(1.0), Inches(2.82), Inches(2.5), Pt(18),
    size=8, color=PINK, font="Calibri")
txb(s3, "Whether Mews defines the platform layer of hospitality before Trybe does.",
    Inches(1.0), Inches(3.0), Inches(5.3), Inches(0.5),
    size=11, bold=False, color=INK, font="Calibri")
txb(s3, "Mews has scale and capital. Trybe has cloud-native architecture and the challenger story. If Mews ships native experience inventory first, Trybe becomes a feature inside Mews's platform.",
    Inches(1.0), Inches(3.45), Inches(5.3), Inches(0.6),
    size=10, color=INK2, font="Calibri")

# Card 02
c2 = glass_box(s3, Inches(6.7), Inches(2.4), Inches(5.65), Inches(1.9))
txb(s3, "02", Inches(6.8), Inches(2.5), Inches(1.5), Pt(30),
    size=22, color=PINK, font="Calibri")
txb(s3, "[WEAKNESS]", Inches(6.8), Inches(2.82), Inches(2.5), Pt(18),
    size=8, color=PINK, font="Calibri")
txb(s3, "Whether the new senior team performs as a unit within 12 months.",
    Inches(6.8), Inches(3.0), Inches(5.3), Inches(0.5),
    size=11, bold=False, color=INK, font="Calibri")
txb(s3, "The senior layer is being assembled in real time. If they integrate and perform fast, everything else is solvable. If they don't, no strategy compensates.",
    Inches(6.8), Inches(3.45), Inches(5.3), Inches(0.6),
    size=10, color=INK2, font="Calibri")

# Three secondary variables
txb(s3, "Three variables that look secondary but aren't",
    Inches(0.9), Inches(4.45), Inches(11.5), Pt(18), size=9, color=INK2, font="Calibri")

sec_cards = [
    ("[Weakness / Threat]", "Brand category-association",
     "A perception problem that takes 18–24 months. Today buyers associate Trybe with 'spa software.' The pivot requires displacing that with 'experience platform.'"),
    ("[Strength / Opportunity]", "The Carden Park story",
     "A 1,500% revenue uplift for a single customer is the most under-leveraged asset Trybe has. The work: turning it into the canonical reference point for the category."),
    ("[Opportunity]", "The chain-wide deal",
     "Winning Marriott, Hilton or Hyatt isn't a single event — it's the moment everything starts compounding: validation, hiring quality, narrative authority, Series B valuation."),
]
sx = Inches(0.9)
for stag, sh, sp in sec_cards:
    glass_box(s3, sx, Inches(4.7), Inches(3.72), Inches(1.6))
    txb(s3, stag, sx + Inches(0.15), Inches(4.78), Inches(3.5), Pt(18),
        size=8, color=MUTED, font="Calibri")
    txb(s3, sh, sx + Inches(0.15), Inches(4.98), Inches(3.5), Pt(20),
        size=11, bold=True, color=INK, font="Calibri")
    txb(s3, sp, sx + Inches(0.15), Inches(5.22), Inches(3.5), Inches(0.9),
        size=10, color=INK2, font="Calibri")
    sx += Inches(3.85)

# Punch line
glass_box(s3, Inches(0.9), Inches(6.45), Inches(11.5), Inches(0.7))
txb(s3, "▶  Two things to get right. Three things to lean into. Everything else is the work, but not the bet.",
    Inches(1.05), Inches(6.57), Inches(10.0), Pt(24),
    size=12, color=INK, font="Calibri")

# ── Slide 4 — Philosophy ───────────────────────────────────────────────────────
s4 = add_slide()
chrome_top(s4, "04 · Performance as the system")

label_txt(s4, "The operating philosophy", Inches(0.9), Inches(0.85))
txb(s4, "Performance is the system.",
    Inches(0.9), Inches(1.2), Inches(11.5), Inches(0.8),
    size=34, color=INK, font="Calibri")
txb(s4, "The pivot is principally an execution question. Execution is principally a people question. The people system must be designed around performance, not process.",
    Inches(0.9), Inches(2.05), Inches(11.5), Inches(0.55),
    size=13, color=INK2, font="Calibri")

# Left card - Problem with traditional HR
lc = glass_box(s4, Inches(0.9), Inches(2.75), Inches(5.65), Inches(4.4))
txb(s4, "The problem with traditional HR",
    Inches(1.05), Inches(2.9), Inches(5.3), Pt(24),
    size=14, bold=True, color=INK, font="Calibri")
txb(s4, "Most People functions treat the lifecycle — recruit, onboard, perform, develop, retain, exit — as the operating model. Each element managed in isolation. Activity confused with progress. Over-engineered for compliance.\n\nThe result: false momentum, no learning, no performance lift.\n\nThe Trybe operating system inverts this. Performance is the core. The HR lifecycle exists to build capability and drive output — not the other way around.",
    Inches(1.05), Inches(3.2), Inches(5.3), Inches(3.5),
    size=12, color=INK2, font="Calibri")

# Right card - Principles
rc = glass_box(s4, Inches(6.7), Inches(2.75), Inches(5.65), Inches(4.4))
txb(s4, "Operating principles",
    Inches(6.85), Inches(2.9), Inches(5.3), Pt(24),
    size=14, bold=True, color=INK, font="Calibri")

principles = [
    "We don't scale with headcount. We scale with performance.",
    "Talent density rises every quarter. Every new hire raises the average, not lowers it.",
    "Every role drives revenue or efficiency. Borrow > Build > Automate > Buy.",
    "Managers are performance enablers — manager-makers, not administrators.",
    "The system carries the standard. The founders stop having to.",
]
py = Inches(3.22)
for i, pr in enumerate(principles):
    txb(s4, f"0{i+1}", Inches(6.85), py, Inches(0.4), Pt(20),
        size=9, color=PINK, font="Calibri")
    txb(s4, pr, Inches(7.2), py, Inches(5.0), Pt(20),
        size=12, color=INK, font="Calibri")
    py += Inches(0.6)

# ── Slide 5 — Operational excellence ──────────────────────────────────────────
s5 = add_slide()
chrome_top(s5, "05 · What I'd measure")

label_txt(s5, "Operational excellence", Inches(0.9), Inches(0.85))
txb(s5, "A performance dashboard, not a sentiment dashboard.",
    Inches(0.9), Inches(1.2), Inches(11.5), Inches(0.8),
    size=30, color=INK, font="Calibri")
txb(s5, "The dashboard tells the founders whether the People system is working, and surfaces leading indicators of the pivot succeeding or stalling.",
    Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.5),
    size=12, color=INK2, font="Calibri")

# Metrics table
bx2 = Inches(0.9)
by2 = Inches(2.6)
bw2 = Inches(8.9)
bh2 = Inches(4.55)
glass_box(s5, bx2, by2, bw2, bh2)

col_ws5 = [Inches(1.4), Inches(3.6), Inches(3.9)]
# Header
hrow = ["Category", "Metric", "Why it matters for the pivot"]
cx5 = bx2 + Inches(0.1)
for hdr, cw in zip(hrow, col_ws5):
    txb(s5, hdr.upper(), cx5, by2 + Inches(0.08), cw, Pt(16),
        size=8, color=MUTED, font="Calibri")
    cx5 += cw

met_rows = [
    ("Performance",    "Team OKR attainment; individual goal attainment; manager effectiveness",
                       "Direct read on whether the company is delivering against the £30m plan"),
    ("Talent density", "% of roles meeting or exceeding Impact Profile standards",
                       "The compounding lever — every hire raises the average or lowers it"),
    ("Hiring",         "Time to hire; quality of hire (6M Impact Profile pass rate); offer acceptance",
                       "Determines whether the senior team integrates and performs"),
    ("Retention",      "Regrettable turnover; engagement pulse top/bottom 3; AMO conditions",
                       "Early signal of senior team health and culture under strain"),
    ("Org health",     "Headcount growth vs revenue per head; cost to serve per property",
                       "Direct read on the sublinear hiring thesis"),
    ("Equity",         "Pay equity by role, gender, ethnicity; promotion distribution",
                       "Protects the no-ego culture as the company scales"),
]
ry5 = by2 + Inches(0.35)
rh5 = Inches(0.64)
col_colors5 = [PINK, INK, INK2]
col_sizes5  = [9, 11, 11]
for cat, metric, why in met_rows:
    cx5 = bx2 + Inches(0.1)
    for txt, cw, col, sz in zip([cat, metric, why], col_ws5, col_colors5, col_sizes5):
        txb(s5, txt, cx5, ry5 + Inches(0.04), cw, rh5 - Inches(0.08),
            size=sz, color=col, font="Calibri")
        cx5 += cw
    ry5 += rh5

# Cadence sidebar
cs = glass_box(s5, Inches(9.95), Inches(2.6), Inches(2.45), Inches(4.55))
txb(s5, "Reporting cadence", Inches(10.1), Inches(2.72), Inches(2.2), Pt(20),
    size=12, bold=True, color=INK, font="Calibri")
cads = [
    "Monthly to Ricky in our 1:1",
    "Quarterly to the full founder team",
    "Three-slide People appendix in the Five Elms board pack: talent density, hiring quality, performance system health",
]
cy = Inches(3.1)
for c in cads:
    txb(s5, f"▸  {c}", Inches(10.1), cy, Inches(2.2), Inches(0.8),
        size=10, color=INK2, font="Calibri")
    cy += Inches(0.7)

# ── Slide 6 — High performance / Impact Profile ────────────────────────────────
s6 = add_slide()
chrome_top(s6, "06 · Impact Profile")

label_txt(s6, "High performance — built not hoped for", Inches(0.9), Inches(0.85))
txb(s6, "The Impact Profile: from task list to success criteria.",
    Inches(0.9), Inches(1.2), Inches(11.5), Inches(0.8),
    size=30, color=INK, font="Calibri")
txb(s6, "Applied to your Customer Onboarding Executive (German Speaker) role — currently live on careers.",
    Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.45),
    size=12, color=INK2, font="Calibri")

# Left blocks
blocks = [
    ("01 · Team alignment",    "Six-month goals, top-down",
     "Skills and capability gap analysis. Team OKRs aligned to the experience platform journey."),
    ("02 · Individual clarity", "Impact Profiles for every role",
     "Metrics, behaviours, Good/Great/Excellent standards. Three-month structured reviews against the Profile."),
    ("03 · Continuous feedback","Weekly 1:1s + quarterly development",
     "Focused on performance and unblocking. Biannual upward manager review."),
    ("04 · Reward",             "50th percentile + performance variable",
     "Equity philosophy that wins senior hires on trajectory. Transparent reward criteria."),
]
blk_y = Inches(2.55)
for n, bh, bp in blocks:
    txb(s6, n, Inches(0.9), blk_y, Inches(3.3), Pt(16),
        size=8, color=PINK, font="Calibri")
    txb(s6, bh, Inches(0.9), blk_y + Inches(0.22), Inches(3.3), Pt(20),
        size=12, bold=True, color=INK, font="Calibri")
    txb(s6, bp, Inches(0.9), blk_y + Inches(0.5), Inches(3.3), Inches(0.5),
        size=10, color=INK2, font="Calibri")
    blk_y += Inches(1.15)

# Right: IP table
tbx = Inches(4.35)
tby = Inches(2.55)
tbw = Inches(8.08)
tbh = Inches(4.6)
glass_box(s6, tbx, tby, tbw, tbh)

ip_headers = ["Job description today", "Impact Profile · what success looks like"]
ip_col_ws  = [Inches(3.0), Inches(5.0)]
cx6 = tbx + Inches(0.1)
for hdr, cw in zip(ip_headers, ip_col_ws):
    col = MUTED if "today" in hdr else PINK
    txb(s6, hdr, cx6, tby + Inches(0.08), cw, Pt(18),
        size=8, color=col, font="Calibri")
    cx6 += cw

ip_rows = [
    ("Guide new customers through onboarding",
     "Time to first booking ≤ 21 days; time to first £10k revenue ≤ 60 days; configuration accuracy ≥ 95%"),
    ("Deliver training and product demos",
     "Training completion ≥ 95%; customer self-sufficiency at go-live"),
    ("Be the main point of contact through handover",
     "Customer health at 90 days ≥ 8/10; handover quality ≥ 9/10; zero unresolved blockers at handover"),
    ("Collaborate cross-functionally",
     "Issue resolution ≤ 48 hours; sales handover satisfaction ≥ 9/10"),
    ("Improve onboarding processes",
     "Process improvements shipped ≥ 1 per quarter; contributions to the configuration library"),
    ("German market: translation and support",
     "KB German coverage ≥ 90% of EN within 6 months; German market NPS ≥ 9"),
]
ry6 = tby + Inches(0.35)
rh6 = Inches(0.68)
for jd, ip in ip_rows:
    txb(s6, jd, tbx + Inches(0.1), ry6 + Inches(0.04),
        ip_col_ws[0] - Inches(0.1), rh6 - Inches(0.06),
        size=10, color=INK2, font="Calibri")
    txb(s6, ip, tbx + Inches(0.1) + ip_col_ws[0], ry6 + Inches(0.04),
        ip_col_ws[1] - Inches(0.1), rh6 - Inches(0.06),
        size=10, color=INK, font="Calibri")
    ry6 += rh6

# ── Slide 7 — Operational gameshare ───────────────────────────────────────────
s7 = add_slide()
chrome_top(s7, "07 · Commercial value")

label_txt(s7, "The operational gameshare", Inches(0.9), Inches(0.85))
txb(s7, "Where People work creates commercial value.",
    Inches(0.9), Inches(1.2), Inches(11.5), Inches(0.8),
    size=34, color=INK, font="Calibri")
txb(s7, "Three operational redesigns that change the unit economics. Each one is a People function intervention with a commercial number attached.",
    Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.5),
    size=12, color=INK2, font="Calibri")

rd_cards = [
    ("01 · Customer onboarding", "Onboarding cost per property",
     "8 onboarders, ~400 properties. Every setup done from scratch. The redesign: three specialised functions (architect, partner, ops), templated configuration library, AI on data migration, menu translation, intake forms, training collateral.",
     "Avoided hiring: ~£400–600k/yr · Concurrent load per person doubles"),
    ("02 · Sales and BD", "Sales productivity",
     "11 people in sales and BD, heavily weighted to BDRs doing manual prospecting. The redesign: AI handles list-building, enrichment, first-touch personalisation, follow-up cadences. Humans handle discovery, demo, negotiation, close.",
     "Avoided hiring: ~£200–300k/yr · Pipeline quality improves"),
    ("03 · Customer success", "Expansion revenue from existing customers",
     "3 people in CX, structured for support tickets rather than account expansion. The redesign: CS becomes a strategic function owning expansion revenue. AI handles tier-one tickets. Humans handle commercial conversations and account growth.",
     "NRR from 100% → 115–120% = £600k–1m+/yr compounding"),
]
rx = Inches(0.9)
for rn, rh, rp, rs in rd_cards:
    glass_box(s7, rx, Inches(2.65), Inches(3.8), Inches(3.45))
    txb(s7, rn, rx + Inches(0.15), Inches(2.78), Inches(3.5), Pt(18),
        size=8, color=PINK, font="Calibri")
    txb(s7, rh, rx + Inches(0.15), Inches(2.98), Inches(3.5), Pt(24),
        size=14, bold=True, color=INK, font="Calibri")
    txb(s7, rp, rx + Inches(0.15), Inches(3.32), Inches(3.5), Inches(1.7),
        size=10, color=INK2, font="Calibri")
    txb(s7, rs, rx + Inches(0.15), Inches(5.52), Inches(3.5), Pt(20),
        size=9, color=PINK, font="Calibri")
    rx += Inches(3.93)

# Headline bar
glass_box(s7, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.85))
txb(s7, "Combined: approximately £1.2–2m per year by year two, compounding from there. The Head of People role costs roughly £150–200k all-in.",
    Inches(1.05), Inches(6.42), Inches(9.0), Inches(0.65),
    size=13, color=INK, font="Calibri")
txb(s7, "Role pays for itself 6–10× before\nperformance, culture, or retention\nwork is counted.",
    Inches(10.2), Inches(6.38), Inches(2.0), Inches(0.7),
    size=9, color=MUTED, align=PP_ALIGN.RIGHT, font="Calibri")

# ── Slide 8 — Talent management ────────────────────────────────────────────────
s8 = add_slide()
chrome_top(s8, "08 · Development, flight risk, succession")

label_txt(s8, "Talent management", Inches(0.9), Inches(0.85))
txb(s8, "One system, three answers.",
    Inches(0.9), Inches(1.2), Inches(11.5), Inches(0.8),
    size=34, color=INK, font="Calibri")
txb(s8, "Development, flight risk and succession are all outputs of the same performance system — not three separate programmes.",
    Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.5),
    size=12, color=INK2, font="Calibri")

quad = [
    ("Development needs", "Surface from the gap", [
        "Gap between current performance and Impact Profile standards surfaces in quarterly development conversations",
        "L&D budget allocated to highest-leverage capability gaps, not generic training",
        "AI tooling and manager coaching prioritised over external programmes",
    ]),
    ("Flight risks", "Three signals", [
        "Engagement pulse decline, manager review patterns, and AMO conditions (Ability, Motivation, Opportunity)",
        "Particular focus on the new senior layer in their first 12 months",
        "Intervene before flight risks become resignations",
    ]),
    ("Succession", "Every senior role covered", [
        "Documented successor candidate for every senior role within 12 months of being filled",
        "Talent reviews twice per year with the founders",
        "Five Elms portfolio network used as both a source and destination",
    ]),
    ("The retention question", "Protect, not fix", [
        "Current retention is unusually strong: 2–3 voluntary leavers, 1 termination",
        "The work isn't fixing retention — it's protecting it through a doubling of headcount and senior team integration",
        "The hardest period is the next 12 months. The system has to prevent the typical deterioration.",
    ]),
]

qx = [Inches(0.9), Inches(6.7), Inches(0.9), Inches(6.7)]
qy = [Inches(2.6), Inches(2.6), Inches(5.0), Inches(5.0)]
qw = Inches(5.65)
qh = Inches(2.25)

for i, (qlabel, qtitle, qbullets) in enumerate(quad):
    glass_box(s8, qx[i], qy[i], qw, qh)
    txb(s8, qlabel.upper(), qx[i] + Inches(0.15), qy[i] + Inches(0.12),
        qw - Inches(0.2), Pt(16), size=8, color=PINK, font="Calibri")
    txb(s8, qtitle, qx[i] + Inches(0.15), qy[i] + Inches(0.35),
        qw - Inches(0.2), Pt(24), size=13, bold=True, color=INK, font="Calibri")
    by8 = qy[i] + Inches(0.7)
    for b in qbullets:
        txb(s8, f"▸  {b}", qx[i] + Inches(0.15), by8,
            qw - Inches(0.25), Inches(0.44),
            size=10, color=INK2, font="Calibri")
        by8 += Inches(0.44)

# ── Slide 9 — Operating rhythms ────────────────────────────────────────────────
s9 = add_slide()
chrome_top(s9, "09 · The cadence")

label_txt(s9, "Operating rhythms", Inches(0.9), Inches(0.85))
txb(s9, "The cadence that sustains performance at scale.",
    Inches(0.9), Inches(1.2), Inches(11.5), Inches(0.8),
    size=34, color=INK, font="Calibri")
txb(s9, "Five rhythms; each designed for a specific failure mode if missing. Rhythm is what makes the system institutional before the company outgrows founder bandwidth.",
    Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.5),
    size=12, color=INK2, font="Calibri")

rhythms = [
    ("Daily / Weekly", [
        "Weekly 1:1 — every manager to every direct report; focused on performance and unblocking",
        "Wednesday team lunch (protect)",
        "Weekly Ricky + Head of People sync for fast decisions",
    ]),
    ("Monthly", [
        "All hands (extend to reinforce the pivot narrative)",
        "Leadership Team meeting (consider fortnightly)",
        "People dashboard review with Ricky",
    ]),
    ("Quarterly", [
        "OKR setting and review: company, team, individual",
        "Performance and development conversations",
        "Talent review with founders",
        "Quarterly party (protect)",
    ]),
    ("Semi-annual", [
        "Formal performance review cycle",
        "Engagement pulse and AMO analysis",
        "Talent strategy refresh",
    ]),
    ("Annual", [
        "Planning aligned to Five Elms board cycle",
        "Compensation review and equity refresh",
        "Values and culture review",
    ]),
]

col_w9 = Inches(2.2)
rx9 = Inches(0.9)
for period, items in rhythms:
    glass_box(s9, rx9, Inches(2.6), col_w9, Inches(3.8))
    txb(s9, period.upper(), rx9 + Inches(0.12), Inches(2.72),
        col_w9 - Inches(0.15), Pt(18), size=8, color=PINK, font="Calibri")
    iy = Inches(3.02)
    for item in items:
        txb(s9, f"▸  {item}", rx9 + Inches(0.12), iy,
            col_w9 - Inches(0.18), Inches(0.65),
            size=10, color=INK2, font="Calibri")
        iy += Inches(0.7)
    rx9 += col_w9 + Inches(0.11)

# Protect note
glass_box(s9, Inches(0.9), Inches(6.52), Inches(11.5), Inches(0.72))
txb(s9, "▶  The Wednesday lunch and the quarterly party are cultural assets, not nice-to-haves. Most People functions let them drift. The work is keeping them deliberate.",
    Inches(1.05), Inches(6.65), Inches(10.5), Inches(0.5),
    size=12, color=INK, font="Calibri")

# ── Slide 10 — Hiring ──────────────────────────────────────────────────────────
s10 = add_slide()
chrome_top(s10, "10 · Discipline at 50th percentile pay")

label_txt(s10, "Hiring", Inches(0.9), Inches(0.85))
txb(s10, "Discipline that wins at 50th percentile pay.",
    Inches(0.9), Inches(1.2), Inches(11.5), Inches(0.8),
    size=32, color=INK, font="Calibri")
txb(s10, "Senior candidates don't apply for jobs. They invest career capital. Our job is to make the ROI obvious and the thesis real.",
    Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.5),
    size=12, color=INK2, font="Calibri")

# Left column: Borrow/Build/Automate/Buy + AI stack
lc10 = glass_box(s10, Inches(0.9), Inches(2.6), Inches(5.65), Inches(2.0))
txb(s10, "The principle — Borrow · Build · Automate · Buy",
    Inches(1.05), Inches(2.72), Inches(5.3), Pt(20),
    size=12, bold=True, color=INK, font="Calibri")
borrow_items = [
    ("Borrow",    "Fractional and consultant talent for specialist work before committing to FTE"),
    ("Build",     "Identify rising talent inside the existing team and promote into stretch roles before going external"),
    ("Automate",  "AI tooling for sourcing, scheduling, screening, reference checks — most manual hiring work compressible by 30–50%"),
    ("Buy",       "External senior hires only when borrowed, built and automated options have been exhausted"),
]
by10 = Inches(3.02)
for k, v in borrow_items:
    txb(s10, k.upper(), Inches(1.05), by10, Inches(0.9), Pt(18),
        size=8, color=PINK, font="Calibri")
    txb(s10, v, Inches(2.0), by10, Inches(4.4), Pt(18),
        size=10, color=INK2, font="Calibri")
    by10 += Inches(0.38)

ai_box = glass_box(s10, Inches(0.9), Inches(4.75), Inches(5.65), Inches(2.4))
txb(s10, "The AI stack I bring in",
    Inches(1.05), Inches(4.87), Inches(5.3), Pt(20),
    size=12, bold=True, color=INK, font="Calibri")
ai_items = [
    ("01", "PDL Sourcing Agent — maps target candidate pools, especially senior commercial talent"),
    ("02", "Interview Transcript Scorer — calibrates assessment across interviewers, removes bottleneck"),
    ("03", "JD Architect — role specs from Impact Profiles in minutes, not days"),
    ("04", "Comp Benchmarker — live competitor JDs, normalised bands, keeps the envelope honest"),
]
ay = Inches(5.17)
for n, desc in ai_items:
    txb(s10, n, Inches(1.05), ay, Inches(0.35), Pt(18), size=8, color=PINK, font="Calibri")
    txb(s10, desc, Inches(1.45), ay, Inches(5.0), Pt(18), size=10, color=INK2, font="Calibri")
    ay += Inches(0.42)
txb(s10, "Not theoretical. Battle-tested. Transferable. Live demo available.",
    Inches(1.05), Inches(6.92), Inches(5.3), Pt(16),
    size=8, color=MUTED, font="Calibri")

# Right column: What senior candidates evaluate
rc10 = glass_box(s10, Inches(6.7), Inches(2.6), Inches(5.65), Inches(4.55))
txb(s10, "What senior candidates are actually evaluating",
    Inches(6.85), Inches(2.72), Inches(5.3), Pt(20),
    size=12, bold=True, color=INK, font="Calibri")
txb(s10, "In 2026, senior candidates think about their next role as an investment of career capital — two to three years of their life, the company is the asset, and they want to know the ROI.",
    Inches(6.85), Inches(3.06), Inches(5.3), Inches(0.7),
    size=10, color=INK2, font="Calibri")

brand_items = [
    ("Purpose with ROI on skills:",
     "Trybe offers building blocks senior operators rarely get — category-defining work, AI-native build, embedded payments, international scale-up."),
    ("Team that expands the network:",
     "Trybe's senior team is small enough that every relationship is real. The Five Elms portfolio extends the network across 70+ companies."),
    ("Challenger story that builds reputation:",
     "'I was at Trybe when they took on Mews and won' is asymmetric reputation upside."),
    ("Material conditions that protect the downside:",
     "Comp framework v1 inside 90 days; progression framework inside 6 months. The thesis has to be real, not staged."),
]
biy = Inches(3.88)
for k, v in brand_items:
    txb(s10, f"▶  {k}", Inches(6.85), biy, Inches(5.3), Pt(18),
        size=10, bold=True, color=INK, font="Calibri")
    txb(s10, v, Inches(6.85), biy + Inches(0.24), Inches(5.3), Inches(0.4),
        size=10, color=INK2, font="Calibri")
    biy += Inches(0.74)

# ── Slide 11 — Leadership and structure ────────────────────────────────────────
s11 = add_slide()
chrome_top(s11, "11 · Team structure")

label_txt(s11, "Leadership and team structure", Inches(0.9), Inches(0.85))
txb(s11, "Right-sized for what the work actually requires.",
    Inches(0.9), Inches(1.2), Inches(11.5), Inches(0.8),
    size=30, color=INK, font="Calibri")
txb(s11, "Solo for the first 90–120 days while diagnosing. Team of two from month 4. Holding at two through year three because the system carries the standard, not the headcount.",
    Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.5),
    size=12, color=INK2, font="Calibri")

years = [
    ("Year 1 · 40–60 people", "Solo → team of two", [
        "Months 1–3 (solo): Listen-first with every senior person. Audit current state. Initial People strategy co-created with founders.",
        "Month 4–6 (second seat): Impact Profiles for 50% of roles. Performance system pilot. Hiring infrastructure. Manager workshops. Comp framework v1.",
        "Second hire isn't inheriting a finished system — co-building it with me from week one.",
    ]),
    ("Year 2 · 60–75 people", "Team of two · deeper systems", [
        "Team stays at two. The work scales because the systems do, not because headcount does.",
        "AI tooling matures across hiring, performance, onboarding, employee experience.",
        "Second seat grows into deeper specialisation — Talent Partner or People Business Partner.",
    ]),
    ("Year 3 · 75–85 people", "Considered third seat only", [
        "Triggers: US team beyond 15–20 people; step-change in hiring volume at Series B.",
        "Without a clear trigger, team stays at two.",
        "AI is what keeps the team small. The second seat is what makes the function effective.",
    ]),
]

yr_w = Inches(3.0)
yr_x = Inches(0.9)
for when, yrh, yrb in years:
    glass_box(s11, yr_x, Inches(2.6), yr_w, Inches(4.55))
    # Pink progress bar top
    prog_w = [Inches(1.0), Inches(2.0), Inches(3.0)][years.index((when, yrh, yrb))]
    rect(s11, yr_x, Inches(2.6), yr_w, Emu(40000), fill=RULE, line_color=RULE)
    rect(s11, yr_x, Inches(2.6), prog_w, Emu(40000), fill=PINK, line_color=PINK)
    txb(s11, when.upper(), yr_x + Inches(0.15), Inches(2.72),
        yr_w - Inches(0.2), Pt(16), size=8, color=PINK, font="Calibri")
    txb(s11, yrh, yr_x + Inches(0.15), Inches(2.98),
        yr_w - Inches(0.2), Pt(22), size=13, bold=True, color=INK, font="Calibri")
    by11 = Inches(3.35)
    for b in yrb:
        txb(s11, f"▸  {b}", yr_x + Inches(0.15), by11,
            yr_w - Inches(0.25), Inches(0.72),
            size=10, color=INK2, font="Calibri")
        by11 += Inches(0.75)
    yr_x += yr_w + Inches(0.2)

# Founders sidebar
fnd = glass_box(s11, Inches(10.05), Inches(2.6), Inches(2.38), Inches(4.55))
txb(s11, "Operating model with founders",
    Inches(10.2), Inches(2.72), Inches(2.1), Pt(22),
    size=11, bold=True, color=INK, font="Calibri")
founders = [
    ("Ricky", "Daily Slack, weekly 1:1, monthly dashboard review. We make decisions together."),
    ("Will",  "Monthly catch-up; engineering and product People needs through structured cadence."),
    ("Steve", "Monthly catch-up; comp, compliance, US payroll, board reporting alignment."),
]
fy = Inches(3.1)
for fname, fdesc in founders:
    txb(s11, fname.upper(), Inches(10.2), fy, Inches(2.1), Pt(16),
        size=9, color=PINK, font="Calibri")
    txb(s11, fdesc, Inches(10.2), fy + Inches(0.22), Inches(2.1), Inches(0.65),
        size=9, color=INK2, font="Calibri")
    fy += Inches(0.98)
txb(s11, "My role is coaching, calibration, accountability — not policing. The Heads own their teams' performance.",
    Inches(10.2), Inches(6.2), Inches(2.1), Inches(0.7),
    size=9, color=MUTED, font="Calibri")

# ── Slide 12 — 90-day plan ─────────────────────────────────────────────────────
s12 = add_slide()
chrome_top(s12, "12 · Diagnose, design, deploy")

label_txt(s12, "First 90 days", Inches(0.9), Inches(0.85))
txb(s12, "Diagnose. Design. Deploy.",
    Inches(0.9), Inches(1.2), Inches(11.5), Inches(0.8),
    size=36, color=INK, font="Calibri")
txb(s12, "Three phases, each one building the foundation for the next.",
    Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.45),
    size=12, color=INK2, font="Calibri")

phases = [
    ("Days 1–30 · Diagnose", "Baseline the truth.", [
        "Listen first. Meet every person at every senior level. Spend time with Ricky, Will, Steve individually.",
        "Audit current state: policies, processes, tools, comp, performance, hiring, retention.",
        "Map senior team integration risks and opportunities.",
        "Share initial People strategy with leadership; co-create with founders.",
    ], "Founders have a clear picture of the People landscape. Senior team integration risks surfaced before they entrench."),
    ("Days 31–90 · Design", "Define the operating system.", [
        "Impact Profiles for 50% of roles, starting with senior team.",
        "Deploy first pulse engagement survey with AMO framing.",
        "Launch performance system pilot with 2–3 functions.",
        "Implement HR infrastructure: ATS, single source of truth, dashboards.",
        "Comp framework v1 drafted with founders. Second seat hired and onboarded.",
    ], "The system has visible early wins. Leadership trusts the cadence. People processes start moving off the founders' plates."),
    ("Month 4–6 · Deploy", "Ship the full system.", [
        "Impact Profiles complete for 100% of roles.",
        "Full performance cycle launched with upward manager feedback.",
        "Comp framework v1 deployed. Progression framework v1 in place.",
        "Hiring system running with AI tooling; time to hire reduced.",
        "People dashboard live to leadership team.",
    ], "Performance is the spine. The system carries the standard. Trybe is operating with the foundation it needs."),
]

ph_w = Inches(3.8)
ph_x = Inches(0.9)
for i, (when, ph, bullets, diff) in enumerate(phases):
    glass_box(s12, ph_x, Inches(2.55), ph_w, Inches(4.3))
    prog_ws12 = [Inches(1.27), Inches(2.53), Inches(3.8)]
    rect(s12, ph_x, Inches(2.55), ph_w, Emu(40000), fill=RULE, line_color=RULE)
    rect(s12, ph_x, Inches(2.55), prog_ws12[i], Emu(40000), fill=PINK, line_color=PINK)
    txb(s12, when.upper(), ph_x + Inches(0.15), Inches(2.67),
        ph_w - Inches(0.2), Pt(16), size=8, color=PINK, font="Calibri")
    txb(s12, ph, ph_x + Inches(0.15), Inches(2.92),
        ph_w - Inches(0.2), Pt(22), size=14, bold=True, color=INK, font="Calibri")
    by12 = Inches(3.28)
    for b in bullets:
        txb(s12, f"▸  {b}", ph_x + Inches(0.15), by12,
            ph_w - Inches(0.25), Inches(0.5),
            size=10, color=INK2, font="Calibri")
        by12 += Inches(0.52)
    txb(s12, diff, ph_x + Inches(0.15), Inches(6.2),
        ph_w - Inches(0.2), Inches(0.52),
        size=10, color=INK, bold=False, font="Calibri")
    ph_x += ph_w + Inches(0.15)

# Punch
glass_box(s12, Inches(0.9), Inches(7.0), Inches(11.5), Inches(0.28))
txb(s12, "▶  By month 6, the People function is no longer a hiring decision Trybe made. It is an operating advantage Trybe has.",
    Inches(1.05), Inches(7.05), Inches(10.5), Pt(20),
    size=12, color=INK, font="Calibri")

# ── Slide 13 — Why me / close ──────────────────────────────────────────────────
s13 = add_slide()
# Aurora-ish tint for cover slide feel
chrome_top(s13, "13 · The close")

label_txt(s13, "Why me", Inches(0.9), Inches(0.85))
txb(s13, "Three reasons, in commercial terms.",
    Inches(0.9), Inches(1.2), Inches(11.5), Inches(0.8),
    size=34, color=INK, font="Calibri")
txb(s13, "The 50th percentile pay philosophy is a test. The work I'm pitching to is built to pass it.",
    Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.5),
    size=13, color=INK2, font="Calibri")

reasons = [
    ("01 · Track record", "I've done this exact shape of work before.",
     "CPO at Kyra at ~60 people. The same problem Trybe has now: scattered managers, founders carrying the standard, no system to scale it. I redesigned the performance system from v1 (95%+ 'meeting expectations' with sub-60% OKR attainment) to v2 (impact factor model with role-specific commercial metrics). That's the spine of how I'd build at Trybe."),
    ("02 · AI tooling", "I've built the stack already.",
     "Not pitching a stack I'd commission. Pitching a stack I've shipped. PDL sourcing agent, interview transcript scoring, JD Architect, comp benchmarker, 1:1 insight extractor. Trybe inherits a working AI-native People operating model from day one, not a roadmap. Happy to demo any of them live in this conversation."),
    ("03 · Commercial lens", "I think commercially because I've been a commercial operator.",
     "CPO at Kyra. Multi-entity experience at Lionbridge and VoxSmart. I've sat in board meetings, owned numbers, defended growth plans. The People function I'd build at Trybe is designed to deliver the Five Elms thesis — not to win HR awards."),
]
rsx = Inches(0.9)
rs_w = Inches(3.78)
for n, rh13, rp in reasons:
    glass_box(s13, rsx, Inches(2.6), rs_w, Inches(3.1))
    txb(s13, n, rsx + Inches(0.15), Inches(2.72),
        rs_w - Inches(0.2), Pt(16), size=8, color=PINK, font="Calibri")
    txb(s13, rh13, rsx + Inches(0.15), Inches(2.95),
        rs_w - Inches(0.2), Pt(22), size=13, bold=True, color=INK, font="Calibri")
    txb(s13, rp, rsx + Inches(0.15), Inches(3.38),
        rs_w - Inches(0.2), Inches(2.0),
        size=10, color=INK2, font="Calibri")
    rsx += rs_w + Inches(0.18)

# Fit grid (4 items)
fit_box = glass_box(s13, Inches(0.9), Inches(5.9), Inches(11.5), Inches(1.2))
fit_items = [
    ("AI-native",     "The JD asked for someone who can scale automation. I've built the tools — not planned them."),
    ("Hands-on",      "No team needing to be onboarded slowly. Notion policies, dashboards, automations — all shipped personally."),
    ("Manager-maker", "You used the phrase. It's how I think about the role. Performance through line managers, not through central HR."),
    ("Commercial",    "I read the JD, the AI hospitality blog, and the Series A announcement. I'm pitching to the commercial role, not the administrative one."),
]
fit_x = Inches(1.0)
fit_w = Inches(2.75)
for k, v in fit_items:
    txb(s13, k.upper(), fit_x, Inches(6.0), fit_w, Pt(16),
        size=8, color=PINK, font="Calibri")
    txb(s13, v, fit_x, Inches(6.2), fit_w, Inches(0.75),
        size=10, color=INK2, font="Calibri")
    fit_x += fit_w + Inches(0.15)

# Contact
txb(s13, "stephaniefadahunsi@outlook.com",
    Inches(0.9), Inches(7.15), Inches(11.5), Pt(20),
    size=10, color=MUTED, align=PP_ALIGN.CENTER, font="Calibri")

# ── Save ──────────────────────────────────────────────────────────────────────
out = "/home/claude/repo/trybe-deck.pptx"
prs.save(out)
print(f"Saved: {out}")
